from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt


class PPTGenerator:

    def add_text_slide(self, prs, title, content):

        slide = prs.slides.add_slide(
            prs.slide_layouts[5]
        )

        slide.shapes.title.text = title

        textbox = slide.shapes.add_textbox(
            Inches(0.6),
            Inches(1.2),
            Inches(8),
            Inches(5)
        )

        textbox.text_frame.word_wrap = True

        textbox.text_frame.text = content

        for p in textbox.text_frame.paragraphs:
            p.font.size = Pt(18)


    def generate(self, project, health, summary, filename):

        output_dir = Path(
            "outputs/presentations"
        )

        output_dir.mkdir(
            parents=True,
            exist_ok=True
        )

        ppt_path = (
            output_dir /
            f"{filename.replace('.xlsx','')}.pptx"
        )


        prs = Presentation()


        project_info = project.get(
            "project",
            {}
        )

        project_name = project_info.get(
            "name",
            "Unknown Project"
        )

        manager = project_info.get(
            "manager",
            "Not Available"
        )


        # Slide 1
        slide = prs.slides.add_slide(
            prs.slide_layouts[0]
        )

        slide.shapes.title.text = (
            "ProjectPulse AI"
        )

        slide.placeholders[1].text = (
            "Executive Project Health Report\n\n"
            f"Project: {project_name}\n"
            f"Health Status: {health.get('overall_health')}\n"
            f"Health Score: {health.get('health_score')}%"
        )


        # Slide 2
        self.add_text_slide(
            prs,
            "Project Health Snapshot",
            f"""
Project:
{project_name}

Project Manager:
{manager}

Overall Health:
{health.get('overall_health')}

Health Score:
{health.get('health_score')}%

Progress:
{health.get('progress_percent')}%

Completed Tasks:
{health.get('completed_tasks')}

In Progress:
{health.get('in_progress_tasks')}

Not Started:
{health.get('not_started_tasks')}
"""
        )


        # Slide 3
        self.add_text_slide(
            prs,
            "Key Risks and Issues",
            f"""
Red Tasks:
{health.get('red_tasks')}

Yellow Tasks:
{health.get('yellow_tasks')}

Overdue Tasks:
{health.get('overdue_tasks')}

Critical Tasks:
{health.get('critical_tasks')}


Executive View:

• Monitor delayed activities
• Resolve blockers
• Track critical milestones
"""
        )


        # Slide 4
        self.add_text_slide(
            prs,
            "Delivery Outlook",
            f"""
Current Progress:
{health.get('progress_percent')}%


Delivery Confidence:

{health.get('overall_health')}


Focus Areas:

• Complete pending activities
• Improve milestone tracking
• Reduce execution risks
"""
        )


        # Slide 5
        self.add_text_slide(
            prs,
            "Executive Recommendations",
            """
1. Review high-risk activities with owners.

2. Prioritize overdue tasks.

3. Improve stakeholder communication.

4. Track milestones weekly.
"""
        )


        # Slide 6
        self.add_text_slide(
            prs,
            "RAG Methodology",
            """
Health classification is calculated using:

• Schedule health
• Task completion progress
• RAG indicators
• Overdue activities
• Critical task identification


Green:
Project is on track.

Amber:
Project requires monitoring.

Red:
Immediate corrective action required.
"""
        )


        # Slide 7
        self.add_text_slide(
            prs,
            "AI Executive Summary",
            summary[:1500]
        )


        prs.save(
            ppt_path
        )

        return str(ppt_path)